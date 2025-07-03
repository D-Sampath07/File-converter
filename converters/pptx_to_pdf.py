from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

def pptx_to_pdf(pptx_path, output_path):
    if not pptx_path.lower().endswith(".pptx"):
        raise ValueError("Input file must be a PPTX file.")

    try:
        prs = Presentation(pptx_path)
        c = canvas.Canvas(output_path, pagesize=A4)
        width, height = A4

        for idx, slide in enumerate(prs.slides):
            if idx > 0:
                c.showPage()  # Start new page for every slide except the first

            c.setFont("Helvetica-Bold", 16)
            c.drawString(40, height - 50, f"Slide {idx + 1}")
            y = height - 80

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    lines = shape.text.strip().splitlines()
                    for line in lines:
                        if y < 40:
                            c.showPage()
                            y = height - 50
                        c.setFont("Helvetica", 12)
                        c.drawString(60, y, line.strip())
                        y -= 18

        c.save()

    except Exception as e:
        raise RuntimeError(f"Failed to convert PPTX to PDF: {e}")

    return output_path
