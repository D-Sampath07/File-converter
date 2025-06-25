import os
import platform

def convert(input_path, output_path, target_format):
    ext = target_format.lower()

    if ext == 'pdf':
        convert_pptx_to_pdf(input_path, output_path)
    elif ext == 'txt':
        pptx_to_txt(input_path, output_path)
    else:
        raise ValueError(f"Unsupported target format for PPTX: {target_format}")

def convert_pptx_to_pdf(input_path, output_path):
    if platform.system() == "Windows":
        import comtypes.client

        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1

        # Make sure input path is absolute and backslash formatted
        input_path = os.path.abspath(input_path)
        output_path = os.path.abspath(output_path)

        presentation = powerpoint.Presentations.Open(input_path)
        presentation.SaveAs(output_path, 32)  # 32 = PDF
        presentation.Close()
        powerpoint.Quit()
    else:
        os.system(f'unoconv -f pdf -o "{output_path}" "{input_path}"')

def pptx_to_txt(input_path, output_path):
    from pptx import Presentation
    prs = Presentation(input_path)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(text)
